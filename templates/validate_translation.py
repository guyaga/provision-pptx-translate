"""
Provision PPTX Translate - Translation Validation Helper
==========================================================
Compare an original PPTX with its translated version to identify
potential issues: empty text boxes, untranslated text, mismatched
slide/element counts.

Usage:
    python validate_translation.py --original input/presentation.pptx --translated translated/presentation_Spanish.pptx
    python validate_translation.py --original input/presentation.pptx --translated translated/presentation_Spanish.pptx --language Spanish
"""

import argparse
import os
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Project directory
# ---------------------------------------------------------------------------
PROJECT_DIR = Path(__file__).resolve().parent.parent

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx package not installed. Run: pip install python-pptx")
    sys.exit(1)


# ===========================================================================
# Text Extraction (simplified for validation)
# ===========================================================================

def extract_all_texts(prs: Presentation) -> list[dict]:
    """
    Extract all visible text from a presentation for validation.
    Returns list of dicts: {slide_idx, shape_name, text, context}
    """
    texts = []

    for slide_idx, slide in enumerate(prs.slides):
        # Slide shapes
        for shape in slide.shapes:
            _extract_from_shape(shape, slide_idx, "slide", texts)

        # Speaker notes
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if shape.has_text_frame:
                    _extract_from_shape(shape, slide_idx, "notes", texts)

    return texts


def _extract_from_shape(shape, slide_idx: int, context: str, texts: list):
    """Recursively extract text from a shape."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append({
                    "slide_idx": slide_idx,
                    "shape_name": shape.name,
                    "text": text,
                    "context": context,
                })

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append({
                            "slide_idx": slide_idx,
                            "shape_name": shape.name,
                            "text": text,
                            "context": f"{context}/table",
                        })

    # Group shapes
    if shape.shape_type is not None and shape.shape_type == 6:
        try:
            for child in shape.shapes:
                _extract_from_shape(child, slide_idx, context, texts)
        except AttributeError:
            pass


# ===========================================================================
# Validation Checks
# ===========================================================================

def validate_translation(
    original_path: str,
    translated_path: str,
    target_language: str | None = None,
) -> list[dict]:
    """
    Validate a translated PPTX against the original.
    Returns a list of issues found.
    """
    issues = []

    print(f"\n{'=' * 60}")
    print(f"  PPTX Translation Validation")
    print(f"  Original:   {original_path}")
    print(f"  Translated: {translated_path}")
    if target_language:
        print(f"  Language:   {target_language}")
    print(f"{'=' * 60}\n")

    # --- Load presentations ---
    try:
        prs_orig = Presentation(original_path)
    except Exception as e:
        issues.append({"severity": "error", "message": f"Cannot open original PPTX: {e}"})
        return issues

    try:
        prs_trans = Presentation(translated_path)
    except Exception as e:
        issues.append({"severity": "error", "message": f"Cannot open translated PPTX: {e}"})
        return issues

    # --- Check 1: Slide count ---
    orig_slides = len(prs_orig.slides)
    trans_slides = len(prs_trans.slides)
    print(f"  Slide count: original={orig_slides}, translated={trans_slides}")

    if orig_slides != trans_slides:
        issues.append({
            "severity": "error",
            "message": f"Slide count mismatch: original has {orig_slides}, translated has {trans_slides}",
        })

    # --- Check 2: Extract and compare text ---
    print("  Extracting text from original...")
    orig_texts = extract_all_texts(prs_orig)
    print(f"  Original text elements: {len(orig_texts)}")

    print("  Extracting text from translated...")
    trans_texts = extract_all_texts(prs_trans)
    print(f"  Translated text elements: {len(trans_texts)}")

    if len(orig_texts) != len(trans_texts):
        issues.append({
            "severity": "warning",
            "message": f"Text element count mismatch: original has {len(orig_texts)}, translated has {len(trans_texts)}",
        })

    # --- Check 3: Empty text boxes ---
    empty_count = 0
    for t in trans_texts:
        if not t["text"].strip():
            empty_count += 1
            issues.append({
                "severity": "warning",
                "slide": t["slide_idx"] + 1,
                "shape": t["shape_name"],
                "context": t["context"],
                "message": f"Empty text box on slide {t['slide_idx'] + 1} in shape '{t['shape_name']}' ({t['context']})",
            })

    if empty_count:
        print(f"  Empty text boxes found: {empty_count}")

    # --- Check 4: Possibly untranslated text ---
    # Compare translated text against original -- if they match exactly,
    # the text may not have been translated (unless it's a preserved term)
    untranslated_count = 0
    min_len = min(len(orig_texts), len(trans_texts))

    for i in range(min_len):
        orig_text = orig_texts[i]["text"]
        trans_text = trans_texts[i]["text"]

        # Skip very short text (likely numbers, abbreviations, or preserved terms)
        if len(orig_text) <= 3:
            continue

        if orig_text == trans_text:
            untranslated_count += 1
            issues.append({
                "severity": "info",
                "slide": orig_texts[i]["slide_idx"] + 1,
                "shape": orig_texts[i]["shape_name"],
                "context": orig_texts[i]["context"],
                "message": (
                    f"Possibly untranslated on slide {orig_texts[i]['slide_idx'] + 1}: "
                    f"'{orig_text[:60]}{'...' if len(orig_text) > 60 else ''}'"
                ),
            })

    if untranslated_count:
        print(f"  Possibly untranslated segments: {untranslated_count}")

    # --- Check 5: Very long translations that might overflow ---
    overflow_warnings = 0
    for i in range(min_len):
        orig_text = orig_texts[i]["text"]
        trans_text = trans_texts[i]["text"]

        if len(orig_text) > 5 and len(trans_text) > len(orig_text) * 1.5:
            overflow_warnings += 1
            issues.append({
                "severity": "info",
                "slide": trans_texts[i]["slide_idx"] + 1,
                "shape": trans_texts[i]["shape_name"],
                "context": trans_texts[i]["context"],
                "message": (
                    f"Translation 50%+ longer on slide {trans_texts[i]['slide_idx'] + 1}: "
                    f"original={len(orig_text)} chars, translated={len(trans_text)} chars"
                ),
            })

    if overflow_warnings:
        print(f"  Potential text overflow warnings: {overflow_warnings}")

    # --- Summary ---
    errors = [i for i in issues if i["severity"] == "error"]
    warnings = [i for i in issues if i["severity"] == "warning"]
    infos = [i for i in issues if i["severity"] == "info"]

    print(f"\n{'=' * 60}")
    print(f"  Validation Results")
    print(f"{'=' * 60}")
    print(f"  Errors:   {len(errors)}")
    print(f"  Warnings: {len(warnings)}")
    print(f"  Info:     {len(infos)}")

    if errors:
        print(f"\n  ERRORS:")
        for issue in errors:
            print(f"    - {issue['message']}")

    if warnings:
        print(f"\n  WARNINGS:")
        for issue in warnings:
            print(f"    - {issue['message']}")

    if infos:
        print(f"\n  INFO:")
        for issue in infos:
            print(f"    - {issue['message']}")

    if not issues:
        print(f"\n  No issues found. Translation looks good!")

    print(f"{'=' * 60}\n")

    return issues


# ===========================================================================
# CLI Entry Point
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Provision PPTX Translate - Validate Translation"
    )
    parser.add_argument(
        "--original", type=str, required=True,
        help="Path to the original PPTX file",
    )
    parser.add_argument(
        "--translated", type=str, required=True,
        help="Path to the translated PPTX file",
    )
    parser.add_argument(
        "--language", type=str, default=None,
        help="Target language (optional, for reporting)",
    )
    args = parser.parse_args()

    # Resolve paths
    original = args.original if os.path.isabs(args.original) else str(PROJECT_DIR / args.original)
    translated = args.translated if os.path.isabs(args.translated) else str(PROJECT_DIR / args.translated)

    if not os.path.exists(original):
        print(f"ERROR: Original file not found: {original}")
        sys.exit(1)

    if not os.path.exists(translated):
        print(f"ERROR: Translated file not found: {translated}")
        sys.exit(1)

    issues = validate_translation(original, translated, target_language=args.language)

    # Exit with error code if there are errors
    errors = [i for i in issues if i["severity"] == "error"]
    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
