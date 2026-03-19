"""
Provision PPTX Translate - Main Translation Script
====================================================
Translate a PowerPoint presentation to another language while preserving
all design, formatting, animations, and layout.

Uses Gemini 3.1 Pro for translation. Replaces text run-by-run to preserve
character-level formatting (font, size, color, bold, italic, underline).

Usage:
    python translate_pptx.py --input presentation.pptx --language Spanish
    python translate_pptx.py --input presentation.pptx --config config.json
    python translate_pptx.py --input presentation.pptx --language Hebrew --config config.json
"""

import argparse
import json
import os
import re
import sys
import time
from pathlib import Path
from copy import deepcopy

# ---------------------------------------------------------------------------
# Project directory: all paths are relative to this script's parent's parent
# ---------------------------------------------------------------------------
PROJECT_DIR = Path(__file__).resolve().parent.parent

# ---------------------------------------------------------------------------
# Optional imports -- fail gracefully with clear messages
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx package not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from google import genai
except ImportError:
    print("ERROR: google-genai package not installed. Run: pip install google-genai")
    sys.exit(1)


# ===========================================================================
# RTL Language Detection
# ===========================================================================

RTL_LANGUAGES = {
    "arabic", "hebrew", "persian", "farsi", "urdu", "pashto", "sindhi",
    "yiddish", "dari",
}


def is_rtl_language(language_name: str) -> bool:
    """Check if the target language is RTL based on its name."""
    return language_name.strip().lower() in RTL_LANGUAGES


def flip_alignment(alignment) -> PP_ALIGN:
    """Flip text alignment for RTL languages (left <-> right)."""
    if alignment == PP_ALIGN.LEFT:
        return PP_ALIGN.RIGHT
    elif alignment == PP_ALIGN.RIGHT:
        return PP_ALIGN.LEFT
    return alignment  # CENTER, JUSTIFY, etc. stay the same


# ===========================================================================
# Text Extraction
# ===========================================================================

def extract_runs_from_paragraph(paragraph) -> list[dict]:
    """
    Extract all text runs from a paragraph, preserving their index
    so we can replace text back into the correct run.
    """
    runs = []
    for run_idx, run in enumerate(paragraph.runs):
        text = run.text
        if text.strip():  # Only include runs with visible text
            runs.append({
                "text": text,
                "run_idx": run_idx,
            })
    return runs


def extract_text_from_shape(shape, slide_idx: int, context: str = "slide") -> list[dict]:
    """
    Recursively extract all text segments from a shape.
    Returns a list of dicts with location info for replacement.
    """
    segments = []

    # --- Text frame (titles, subtitles, body text, content placeholders) ---
    if shape.has_text_frame:
        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            full_para_text = paragraph.text.strip()
            if not full_para_text:
                continue

            # Collect all runs for this paragraph
            runs = extract_runs_from_paragraph(paragraph)
            if runs:
                segments.append({
                    "slide_idx": slide_idx,
                    "context": context,
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "para_idx": para_idx,
                    "full_text": full_para_text,
                    "runs": runs,
                    "type": "text_frame",
                })

    # --- Table ---
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                    full_para_text = paragraph.text.strip()
                    if not full_para_text:
                        continue
                    runs = extract_runs_from_paragraph(paragraph)
                    if runs:
                        segments.append({
                            "slide_idx": slide_idx,
                            "context": context,
                            "shape_id": shape.shape_id,
                            "shape_name": shape.name,
                            "para_idx": para_idx,
                            "row_idx": row_idx,
                            "col_idx": col_idx,
                            "full_text": full_para_text,
                            "runs": runs,
                            "type": "table_cell",
                        })

    # --- Group shapes (recursive) ---
    if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        try:
            for child_shape in shape.shapes:
                child_segments = extract_text_from_shape(child_shape, slide_idx, context)
                segments.extend(child_segments)
        except AttributeError:
            pass  # Some group shapes may not expose .shapes

    return segments


def extract_all_text(prs: Presentation, translate_notes: bool = True) -> list[dict]:
    """
    Extract all text from a presentation: slide content and speaker notes.
    Returns a flat list of text segments with location metadata.
    """
    all_segments = []

    for slide_idx, slide in enumerate(prs.slides):
        print(f"  Extracting text from slide {slide_idx + 1}/{len(prs.slides)}...")

        # --- Slide shapes ---
        for shape in slide.shapes:
            segments = extract_text_from_shape(shape, slide_idx, context="slide")
            all_segments.extend(segments)

        # --- Speaker notes ---
        if translate_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                # Skip the slide image placeholder in notes
                if shape.has_text_frame:
                    segments = extract_text_from_shape(shape, slide_idx, context="notes")
                    all_segments.extend(segments)

    print(f"  Total text segments extracted: {len(all_segments)}")
    return all_segments


# ===========================================================================
# Translation with Gemini
# ===========================================================================

def build_translation_prompt(
    texts: list[str],
    target_language: str,
    preserve_terms: list[str],
) -> str:
    """Build the Gemini prompt for translating a batch of text segments."""
    terms_str = ", ".join(f'"{t}"' for t in preserve_terms)

    # Build numbered text list
    numbered_texts = []
    for i, text in enumerate(texts):
        numbered_texts.append(f"{i}: {text}")
    texts_block = "\n".join(numbered_texts)

    prompt = f"""You are a professional translator. Translate the following numbered text segments from their original language into {target_language}.

RULES:
1. Translate each segment independently but maintain consistent terminology.
2. Do NOT translate these product/brand terms -- keep them exactly as-is: {terms_str}
3. Keep the same tone and formality as the original.
4. For technical terms, use the standard localized equivalent in {target_language}.
5. Preserve any numbers, units, and formatting markers.
6. If a segment is ONLY a product name or brand term from the preserve list, return it unchanged.

TEXT SEGMENTS:
{texts_block}

Return ONLY a JSON object mapping the segment number (as string) to the translated text.
Example format: {{"0": "translated text 0", "1": "translated text 1", ...}}
Return raw JSON only, no markdown fences, no explanation."""

    return prompt


def translate_batch(
    texts: list[str],
    target_language: str,
    preserve_terms: list[str],
    gemini_model: str,
    client,
) -> dict[str, str]:
    """
    Send a batch of text segments to Gemini for translation.
    Returns a dict mapping index (as string) to translated text.
    """
    prompt = build_translation_prompt(texts, target_language, preserve_terms)

    response = client.models.generate_content(
        model=gemini_model,
        contents=prompt,
    )

    response_text = response.text.strip()
    # Remove markdown code fences if present
    if response_text.startswith("```"):
        response_text = re.sub(r"^```(?:json)?\s*", "", response_text)
        response_text = re.sub(r"\s*```$", "", response_text)

    translations = json.loads(response_text)
    return translations


def translate_all_segments(
    segments: list[dict],
    target_language: str,
    preserve_terms: list[str],
    gemini_model: str,
    batch_size: int = 50,
) -> list[dict]:
    """
    Translate all extracted text segments using Gemini in batches.
    Returns the segments list with 'translated_text' added to each.
    """
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        print("ERROR: GEMINI_API_KEY environment variable not set.")
        sys.exit(1)

    client = genai.Client(api_key=api_key)

    # Collect unique texts to avoid translating duplicates
    unique_texts = []
    text_to_indices = {}  # maps full_text -> list of segment indices
    for idx, seg in enumerate(segments):
        text = seg["full_text"]
        if text not in text_to_indices:
            text_to_indices[text] = []
            unique_texts.append(text)
        text_to_indices[text].append(idx)

    print(f"  Unique text segments to translate: {len(unique_texts)}")

    # Translate in batches
    all_translations = {}  # maps original text -> translated text
    for batch_start in range(0, len(unique_texts), batch_size):
        batch = unique_texts[batch_start:batch_start + batch_size]
        batch_num = batch_start // batch_size + 1
        total_batches = (len(unique_texts) + batch_size - 1) // batch_size
        print(f"  Translating batch {batch_num}/{total_batches} ({len(batch)} segments)...")

        try:
            translations = translate_batch(
                batch, target_language, preserve_terms, gemini_model, client
            )

            # Map translations back to original texts
            for i, text in enumerate(batch):
                idx_str = str(i)
                if idx_str in translations:
                    all_translations[text] = translations[idx_str]
                else:
                    print(f"  WARNING: No translation returned for segment {batch_start + i}, keeping original.")
                    all_translations[text] = text

        except json.JSONDecodeError as e:
            print(f"  ERROR: Failed to parse Gemini response as JSON: {e}")
            print(f"  Keeping original text for this batch.")
            for text in batch:
                all_translations[text] = text

        except Exception as e:
            print(f"  ERROR: Translation batch failed: {e}")
            print(f"  Keeping original text for this batch.")
            for text in batch:
                all_translations[text] = text

        # Brief pause between batches to respect rate limits
        if batch_start + batch_size < len(unique_texts):
            time.sleep(1)

    # Apply translations to all segments
    for seg in segments:
        seg["translated_text"] = all_translations.get(seg["full_text"], seg["full_text"])

    translated_count = sum(1 for s in segments if s["translated_text"] != s["full_text"])
    print(f"  Successfully translated {translated_count}/{len(segments)} segments.")
    return segments


# ===========================================================================
# Text Replacement in PPTX
# ===========================================================================

def distribute_translation_to_runs(original_runs: list[dict], translated_text: str) -> list[str]:
    """
    Distribute translated text across the original runs.

    Strategy:
    - If there is only one run, the entire translated text goes into it.
    - If there are multiple runs, we try to split the translated text
      proportionally based on the original run lengths.
    - This preserves per-run formatting as closely as possible.
    """
    if len(original_runs) == 1:
        return [translated_text]

    # Calculate proportional lengths
    original_lengths = [len(r["text"]) for r in original_runs]
    total_original_len = sum(original_lengths)

    if total_original_len == 0:
        return [translated_text] + [""] * (len(original_runs) - 1)

    # Split translated text by words for cleaner distribution
    words = translated_text.split()
    if not words:
        return [""] * len(original_runs)

    # Distribute words proportionally across runs
    result = []
    word_idx = 0
    for run_i, orig_len in enumerate(original_lengths):
        if run_i == len(original_lengths) - 1:
            # Last run gets all remaining words
            run_words = words[word_idx:]
        else:
            # Proportional word count for this run
            proportion = orig_len / total_original_len
            word_count = max(1, round(proportion * len(words)))
            run_words = words[word_idx:word_idx + word_count]
            word_idx += word_count

        result.append(" ".join(run_words))

    return result


def replace_text_in_paragraph(paragraph, original_runs: list[dict], translated_texts: list[str]):
    """
    Replace text in a paragraph's runs with translated text.
    Preserves all formatting (font, size, color, bold, italic, underline).
    """
    for run_info, new_text in zip(original_runs, translated_texts):
        run_idx = run_info["run_idx"]
        if run_idx < len(paragraph.runs):
            paragraph.runs[run_idx].text = new_text


def apply_translations_to_pptx(
    prs: Presentation,
    segments: list[dict],
    rtl: bool = False,
):
    """
    Apply translated text back into the PPTX, preserving all formatting.
    """
    slides = list(prs.slides)

    for seg in segments:
        slide_idx = seg["slide_idx"]
        translated = seg.get("translated_text", seg["full_text"])

        # Skip if translation is same as original (nothing to do)
        if translated == seg["full_text"]:
            continue

        # Determine which shape collection to search
        if seg["context"] == "notes":
            if not slides[slide_idx].has_notes_slide:
                continue
            shapes = slides[slide_idx].notes_slide.shapes
        else:
            shapes = slides[slide_idx].shapes

        # Find the target shape
        target_shape = None
        for shape in shapes:
            if shape.shape_id == seg["shape_id"]:
                target_shape = shape
                break

        # For group shapes, we may need to search recursively
        if target_shape is None:
            target_shape = _find_shape_recursive(shapes, seg["shape_id"])

        if target_shape is None:
            print(f"  WARNING: Could not find shape {seg['shape_name']} (ID={seg['shape_id']}) "
                  f"on slide {slide_idx + 1}. Skipping.")
            continue

        # --- Replace text based on segment type ---
        if seg["type"] == "table_cell":
            _replace_in_table_cell(target_shape, seg, translated, rtl)
        elif seg["type"] == "text_frame":
            _replace_in_text_frame(target_shape, seg, translated, rtl)


def _find_shape_recursive(shapes, shape_id):
    """Recursively search for a shape by ID within group shapes."""
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        # Check group shapes
        if shape.shape_type is not None and shape.shape_type == 6:
            try:
                found = _find_shape_recursive(shape.shapes, shape_id)
                if found:
                    return found
            except AttributeError:
                pass
    return None


def _replace_in_text_frame(shape, seg: dict, translated: str, rtl: bool):
    """Replace text in a text frame paragraph."""
    para_idx = seg["para_idx"]
    if para_idx >= len(shape.text_frame.paragraphs):
        return

    paragraph = shape.text_frame.paragraphs[para_idx]
    original_runs = seg["runs"]
    translated_texts = distribute_translation_to_runs(original_runs, translated)
    replace_text_in_paragraph(paragraph, original_runs, translated_texts)

    # Flip alignment for RTL languages
    if rtl and paragraph.alignment is not None:
        paragraph.alignment = flip_alignment(paragraph.alignment)


def _replace_in_table_cell(shape, seg: dict, translated: str, rtl: bool):
    """Replace text in a table cell paragraph."""
    row_idx = seg["row_idx"]
    col_idx = seg["col_idx"]
    para_idx = seg["para_idx"]

    table = shape.table
    if row_idx >= len(table.rows):
        return
    cell = table.rows[row_idx].cells[col_idx]
    if para_idx >= len(cell.text_frame.paragraphs):
        return

    paragraph = cell.text_frame.paragraphs[para_idx]
    original_runs = seg["runs"]
    translated_texts = distribute_translation_to_runs(original_runs, translated)
    replace_text_in_paragraph(paragraph, original_runs, translated_texts)

    if rtl and paragraph.alignment is not None:
        paragraph.alignment = flip_alignment(paragraph.alignment)


# ===========================================================================
# Main Pipeline
# ===========================================================================

def load_config(config_path: str) -> dict:
    """Load configuration from JSON file."""
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def translate_pptx(
    input_path: str,
    target_language: str,
    config: dict,
    output_path: str | None = None,
) -> str:
    """
    Main function: translate a PPTX file to the target language.
    Returns the path to the translated PPTX.
    """
    gemini_model = config.get("gemini_model", "gemini-3.1-pro")
    preserve_terms = config.get("preserve_terms", [])
    translate_notes = config.get("translate_speaker_notes", True)
    output_suffix = config.get("output_suffix", True)
    output_dir = config.get("output_dir", "translated")

    # Resolve output path
    if output_path is None:
        input_p = Path(input_path)
        out_dir = PROJECT_DIR / output_dir
        os.makedirs(out_dir, exist_ok=True)
        if output_suffix:
            output_path = str(out_dir / f"{input_p.stem}_{target_language}{input_p.suffix}")
        else:
            output_path = str(out_dir / input_p.name)

    rtl = is_rtl_language(target_language)

    print(f"\n{'=' * 60}")
    print(f"  PPTX Translation Pipeline")
    print(f"  Input:    {input_path}")
    print(f"  Language: {target_language}")
    print(f"  RTL:      {rtl}")
    print(f"  Output:   {output_path}")
    print(f"{'=' * 60}")

    # --- Step 1: Open PPTX ---
    print("\n=== Step 1: Opening PPTX ===")
    prs = Presentation(input_path)
    print(f"  Loaded presentation with {len(prs.slides)} slides.")

    # --- Step 2: Extract all text ---
    print("\n=== Step 2: Extracting Text ===")
    segments = extract_all_text(prs, translate_notes=translate_notes)

    if not segments:
        print("  No text found to translate. Saving copy as-is.")
        prs.save(output_path)
        return output_path

    # --- Step 3: Translate with Gemini ---
    print("\n=== Step 3: Translating with Gemini ===")
    segments = translate_all_segments(
        segments,
        target_language=target_language,
        preserve_terms=preserve_terms,
        gemini_model=gemini_model,
        batch_size=50,
    )

    # --- Step 4: Apply translations back to PPTX ---
    print("\n=== Step 4: Applying Translations ===")
    apply_translations_to_pptx(prs, segments, rtl=rtl)

    # --- Step 5: Save translated PPTX ---
    print("\n=== Step 5: Saving Translated PPTX ===")
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"  Saved translated PPTX: {output_path}")

    print(f"\n{'=' * 60}")
    print(f"  Translation complete!")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Text segments translated: {len(segments)}")
    print(f"  Output: {output_path}")
    print(f"{'=' * 60}\n")

    return output_path


# ===========================================================================
# CLI Entry Point
# ===========================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Provision PPTX Translate - Translate PowerPoint presentations"
    )
    parser.add_argument(
        "--input", type=str, required=True,
        help="Path to the source PPTX file",
    )
    parser.add_argument(
        "--language", type=str, default=None,
        help="Target language (e.g., Spanish, Italian, Hebrew). Overrides config.",
    )
    parser.add_argument(
        "--config", type=str, default=None,
        help="Path to config.json file",
    )
    parser.add_argument(
        "--output", type=str, default=None,
        help="Path for the translated output PPTX (optional)",
    )
    args = parser.parse_args()

    # Load config
    if args.config:
        config_path = args.config if os.path.isabs(args.config) else str(PROJECT_DIR / args.config)
        config = load_config(config_path)
    else:
        config = {}

    # Determine target language
    target_language = args.language or config.get("target_language", "Spanish")

    # Resolve input path
    input_path = args.input if os.path.isabs(args.input) else str(PROJECT_DIR / args.input)

    if not os.path.exists(input_path):
        print(f"ERROR: Input file not found: {input_path}")
        sys.exit(1)

    translate_pptx(input_path, target_language, config, output_path=args.output)


if __name__ == "__main__":
    main()
